import pandas as pd
from pptx import Presentation
import re
import requests
import imageio
import os
import win32com.client
from tqdm import tqdm
import PySimpleGUI as sg

def colNameToNum(name):
    pow_ = 1
    colNum = 0
    for letter in name[::-1]:
            colNum += (int(letter, 36) -9) * pow_
            pow_ *= 26
    return colNum

def duplicate_slide(pptx_input_file, nbr_duplicates):
    ppt_instance = win32com.client.Dispatch('PowerPoint.Application')
    #open the powerpoint presentation headless in background
    read_only = True
    has_title = False
    window    = False
    prs = ppt_instance.Presentations.open(os.path.abspath(pptx_input_file), read_only, has_title,window)
    nr_slide = 1
    insert_index = 1
    i=0
    print('-------------------Duplicating Slides-----------------')
    pbar = tqdm(total = nbr_duplicates)
    while i<nbr_duplicates-1:
        try:
            prs.Slides(nr_slide).Copy()
            prs.Slides.Paste(Index=insert_index)
            insert_index += 1
            i+=1
            pbar.update(1)
        except:
            pass

    prs.SaveAs(os.getcwd()+'\\test.pptx')
    prs.Close()

    #kills ppt_instance
    ppt_instance.Quit()
    del ppt_instance
    
def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text
    
    
def main():
    
    sg.theme("DarkGreen2")
    layout = [[sg.T("")], [sg.Text("Choose the Excel file: "), sg.Input(), sg.FileBrowse(key="-IN-")],[sg.Button("Submit")]]

    ###Building Window
    window = sg.Window('Excel file', layout, size=(600,150))

    while True:
        event, values = window.read()
        if event == sg.WIN_CLOSED or event=="Exit":
            break
        elif event == "Submit":
            excel_input_file = values["-IN-"]
            window.close()
            break


    sg.theme("LightBrown6")
    layout = [[sg.T("")], [sg.Text("Choose the PPTX template file: "), sg.Input(), sg.FileBrowse(key="-IN-")],[sg.Button("Submit")]]

    ###Building Window
    window = sg.Window('Output Folder', layout, size=(600,150))

    while True:
        event, values = window.read()
        if event == sg.WIN_CLOSED or event=="Exit":
            break
        elif event == "Submit":
            pptx_input_file = values["-IN-"]
            window.close()
            break
            
    df = pd.read_excel(excel_input_file)
    df = df[df[df.columns[0]].notna()]
    
    duplicate_slide(pptx_input_file, df.shape[0])
    prs = Presentation(os.getcwd()+'\\test.pptx')
    print('----------Populating Slides-----------------')
    matching_text = {}
    matching_pictures = {}
    for j in range(len(prs.slides[0].shapes)):
        shape = prs.slides[0].shapes[j]
        for k in range(len(shape.text_frame.paragraphs)):
            p = shape.text_frame.paragraphs[k]
            if p.text == "":
                continue
            if bool(re.match("\[Excel:.*\]", p.text)):
                matching_text[(j,k)] = colNameToNum(shape.text.strip()[7:-1]) - 1
        if bool(re.match("\[GetPhoto\(Excel:.*\)\]", shape.text)):
            matching_pictures[j] = colNameToNum(shape.text.strip()[16:-2]) - 1



    for i in tqdm(range(df.shape[0])):
        for shape_id, paragraph_nbr in matching_text.keys():
            try:
                text = df.iloc[i,matching_text[(shape_id,paragraph_nbr)]]
            except:
                text = ""
            if pd.isna(text):
                text = ""
            else:
                text = str(text)
            replace_paragraph_text_retaining_initial_formatting(prs.slides[i].shapes[shape_id].text_frame.paragraphs[paragraph_nbr], text);
        for shape_id in matching_pictures:
            shape = prs.slides[i].shapes[shape_id]
            left = shape.left
            width = shape.width
            height = shape.height
            top = shape.top
            try:
                img_data = requests.get(df.iloc[i,matching_pictures[shape_id]]).content
                with open(f'current_image.jpg', 'wb') as handler:
                    handler.write(img_data)
                pic   = prs.slides[i].shapes.add_picture('current_image.jpg', left, top-100000, width, height*1.05)
                os.remove("current_image.jpg") 
            except:
                shape.text = ""
                pass   
    
      
    output_file = "output.pptx"
    prs.save(output_file)
    os.remove(os.getcwd()+'\\test.pptx')
    
if __name__ == '__main__':
    main()